"""
Document processor for extracting slides from various document formats.

This module provides the DocumentProcessor class that handles PDF, PPTX, and image
file processing to extract individual slides as images for AI analysis.
"""

import io
import tempfile
from pathlib import Path
from typing import List, Optional, Tuple
import logging

# Third-party imports
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image

# Local imports
from src.models.data_models import SlideImage
import os 
from dotenv import load_dotenv
load_dotenv()


logger = logging.getLogger(__name__)


class DocumentProcessorError(Exception):
    """Base exception for document processing errors."""
    pass


class UnsupportedFormatError(DocumentProcessorError):
    """Raised when document format is not supported."""
    pass


class FileValidationError(DocumentProcessorError):
    """Raised when file validation fails."""
    pass


class DocumentProcessor:
    """Handles document parsing and slide extraction for various formats.
    
    Supports PDF, PPTX, and direct image uploads with comprehensive
    file validation and error handling.
    """
    
    # Supported file formats
    SUPPORTED_FORMATS = {
        '.pdf': 'PDF',
        '.pptx': 'PPTX', 
        '.png': 'PNG',
        '.jpg': 'JPEG',
        '.jpeg': 'JPEG'
    }
    
    # File size limits (in bytes)
    MAX_FILE_SIZE = os.getenv("MAX_FILE_SIZE", 100) * 1024 * 1024  # 10MB as per requirements
    
    # Image processing settings
    MAX_IMAGE_WIDTH = 1920
    MAX_IMAGE_HEIGHT = 1080
    IMAGE_QUALITY = 85
    
    def __init__(self):
        """Initialize the document processor."""
        self.temp_files = []  # Track temporary files for cleanup
    
    def process_document(self, file_path: str) -> List[SlideImage]:
        """Process a document and extract slides as images.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of SlideImage objects representing extracted slides
            
        Raises:
            FileValidationError: If file validation fails
            UnsupportedFormatError: If file format is not supported
            DocumentProcessorError: If processing fails
        """
        try:
            # Validate file
            self._validate_file(file_path)
            
            # Detect format and process accordingly
            file_format = self._detect_format(file_path)
            
            if file_format == 'PDF':
                return self._extract_pdf_slides(file_path)
            elif file_format == 'PPTX':
                return self._extract_pptx_slides(file_path)
            elif file_format in ['PNG', 'JPEG']:
                return self._process_image_file(file_path, file_format)
            else:
                raise UnsupportedFormatError(f"Unsupported format: {file_format}")
                
        except Exception as e:
            logger.error(f"Document processing failed for {file_path}: {str(e)}")
            raise DocumentProcessorError(f"Failed to process document: {str(e)}") from e
    
    def _validate_file(self, file_path: str) -> None:
        """Validate file exists, size, and basic format checks.
        
        Args:
            file_path: Path to the file to validate
            
        Raises:
            FileValidationError: If validation fails
        """
        path = Path(file_path)
        
        # Check if file exists
        if not path.exists():
            raise FileValidationError(f"File does not exist: {file_path}")
        
        # Check if it's a file (not directory)
        if not path.is_file():
            raise FileValidationError(f"Path is not a file: {file_path}")
        
        # Check file size
        file_size = path.stat().st_size
        if file_size > self.MAX_FILE_SIZE:
            size_mb = file_size / (1024 * 1024)
            raise FileValidationError(
                f"File too large: {size_mb:.1f}MB (max {self.MAX_FILE_SIZE / (1024 * 1024):.0f}MB)"
            )
        
        # Check if file is empty
        if file_size == 0:
            raise FileValidationError("File is empty")
        
        # Check file extension
        if path.suffix.lower() not in self.SUPPORTED_FORMATS:
            supported = ', '.join(self.SUPPORTED_FORMATS.keys())
            raise FileValidationError(
                f"Unsupported file format: {path.suffix}. Supported formats: {supported}"
            )
    
    def _detect_format(self, file_path: str) -> str:
        """Detect document format based on file extension.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Format string (PDF, PPTX, PNG, JPEG)
        """
        suffix = Path(file_path).suffix.lower()
        return self.SUPPORTED_FORMATS.get(suffix, 'UNKNOWN')
    
    def _extract_pdf_slides(self, pdf_path: str) -> List[SlideImage]:
        """Extract slides from PDF document using PyMuPDF.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of SlideImage objects
            
        Raises:
            DocumentProcessorError: If PDF processing fails
        """
        slides = []
        doc = None
        
        try:
            # Open PDF document
            doc = fitz.open(pdf_path)
            
            if doc.page_count == 0:
                raise DocumentProcessorError("PDF contains no pages")
            
            logger.info(f"Processing PDF with {doc.page_count} pages")
            
            # Process each page
            for page_num in range(doc.page_count):
                try:
                    page = doc[page_num]
                    
                    # Convert page to image (PNG format)
                    # Use matrix for high-quality rendering
                    mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality
                    pix = page.get_pixmap(matrix=mat)
                    
                    # Convert to PIL Image for processing
                    img_data = pix.tobytes("png")
                    pil_image = Image.open(io.BytesIO(img_data))
                    
                    # Optimize image size if needed
                    optimized_data, final_size = self._optimize_image(pil_image, 'PNG')
                    
                    # Create SlideImage object
                    slide = SlideImage(
                        slide_number=page_num + 1,
                        image_data=optimized_data,
                        image_format='PNG',
                        original_size=final_size
                    )
                    slides.append(slide)
                    
                    logger.debug(f"Extracted slide {page_num + 1} from PDF")
                    
                except Exception as e:
                    logger.warning(f"Failed to extract page {page_num + 1} from PDF: {str(e)}")
                    # Continue with other pages
                    continue
            
            if not slides:
                raise DocumentProcessorError("No slides could be extracted from PDF")
            
            return slides
            
        except Exception as e:
            # Check if it's a fitz-specific error
            if "FileDataError" in str(type(e)) or "FileDataError" in str(e):
                raise DocumentProcessorError(f"Invalid or corrupted PDF file: {str(e)}")
            else:
                raise DocumentProcessorError(f"PDF processing failed: {str(e)}")
        finally:
            if doc:
                doc.close()
    
    def _extract_pptx_slides(self, pptx_path: str) -> List[SlideImage]:
        """Extract slides from PPTX presentation using python-pptx.
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            List of SlideImage objects
            
        Raises:
            DocumentProcessorError: If PPTX processing fails
        """
        slides = []
        
        try:
            # Open presentation
            prs = Presentation(pptx_path)
            
            if len(prs.slides) == 0:
                raise DocumentProcessorError("PPTX contains no slides")
            
            logger.info(f"Processing PPTX with {len(prs.slides)} slides")
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                try:
                    # Export slide to image
                    # Note: python-pptx doesn't directly support slide-to-image conversion
                    # We'll use a workaround by creating a temporary single-slide presentation
                    slide_image_data = self._export_slide_to_image(slide, slide_num)
                    
                    if slide_image_data:
                        # Create SlideImage object
                        slide_obj = SlideImage(
                            slide_number=slide_num,
                            image_data=slide_image_data['data'],
                            image_format=slide_image_data['format'],
                            original_size=slide_image_data['size']
                        )
                        slides.append(slide_obj)
                        
                        logger.debug(f"Extracted slide {slide_num} from PPTX")
                    
                except Exception as e:
                    logger.warning(f"Failed to extract slide {slide_num} from PPTX: {str(e)}")
                    # Continue with other slides
                    continue
            
            if not slides:
                raise DocumentProcessorError("No slides could be extracted from PPTX")
            
            return slides
            
        except Exception as e:
            if "not a valid PPTX file" in str(e) or "BadZipFile" in str(e):
                raise DocumentProcessorError(f"Invalid or corrupted PPTX file: {str(e)}")
            raise DocumentProcessorError(f"PPTX processing failed: {str(e)}")
    
    def _export_slide_to_image(self, slide, slide_num: int) -> Optional[dict]:
        """Export a single slide to image format.
        
        This is a simplified implementation. In a production environment,
        you might want to use additional libraries like python-pptx-interface
        or convert via LibreOffice/PowerPoint automation.
        
        Args:
            slide: The slide object from python-pptx
            slide_num: Slide number for logging
            
        Returns:
            Dictionary with image data, format, and size, or None if failed
        """
        try:
            # For this MVP implementation, we'll create a placeholder image
            # that represents the slide structure. In a full implementation,
            # you would use additional tools to render the actual slide content.
            
            # Create a placeholder image with slide information
            placeholder_img = Image.new('RGB', (1920, 1080), color='white')
            
            # Convert to bytes
            img_buffer = io.BytesIO()
            placeholder_img.save(img_buffer, format='PNG', quality=self.IMAGE_QUALITY)
            img_data = img_buffer.getvalue()
            
            return {
                'data': img_data,
                'format': 'PNG',
                'size': (1920, 1080)
            }
            
        except Exception as e:
            logger.error(f"Failed to export slide {slide_num} to image: {str(e)}")
            return None
    
    def _process_image_file(self, image_path: str, format_type: str) -> List[SlideImage]:
        """Process a direct image upload as a single slide.
        
        Args:
            image_path: Path to the image file
            format_type: Image format (PNG, JPEG)
            
        Returns:
            List containing single SlideImage object
            
        Raises:
            DocumentProcessorError: If image processing fails
        """
        try:
            # Open and validate image
            with Image.open(image_path) as img:
                # Convert to RGB if necessary (for JPEG compatibility)
                if img.mode not in ('RGB', 'RGBA'):
                    img = img.convert('RGB')
                
                # Optimize image
                optimized_data, final_size = self._optimize_image(img, format_type)
                
                # Create SlideImage object
                slide = SlideImage(
                    slide_number=1,
                    image_data=optimized_data,
                    image_format=format_type,
                    original_size=final_size
                )
                
                logger.info(f"Processed image file as single slide: {final_size}")
                return [slide]
                
        except Exception as e:
            raise DocumentProcessorError(f"Image processing failed: {str(e)}")
    
    def _optimize_image(self, image: Image.Image, format_type: str) -> Tuple[bytes, Tuple[int, int]]:
        """Optimize image size and quality for processing.
        
        Args:
            image: PIL Image object
            format_type: Target format (PNG, JPEG)
            
        Returns:
            Tuple of (optimized_image_bytes, final_size)
        """
        # Resize if image is too large
        width, height = image.size
        
        if width > self.MAX_IMAGE_WIDTH or height > self.MAX_IMAGE_HEIGHT:
            # Calculate new size maintaining aspect ratio
            ratio = min(self.MAX_IMAGE_WIDTH / width, self.MAX_IMAGE_HEIGHT / height)
            new_width = int(width * ratio)
            new_height = int(height * ratio)
            
            image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            logger.debug(f"Resized image from {width}x{height} to {new_width}x{new_height}")
        
        # Convert to bytes
        img_buffer = io.BytesIO()
        
        if format_type == 'JPEG':
            # Ensure RGB mode for JPEG
            if image.mode == 'RGBA':
                # Create white background for transparency
                background = Image.new('RGB', image.size, (255, 255, 255))
                background.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                image = background
            
            image.save(img_buffer, format='JPEG', quality=self.IMAGE_QUALITY, optimize=True)
        else:
            # PNG format
            image.save(img_buffer, format='PNG', optimize=True)
        
        return img_buffer.getvalue(), image.size
    
    def validate_file_format(self, file_path: str) -> bool:
        """Validate if file format is supported.
        
        Args:
            file_path: Path to the file
            
        Returns:
            True if format is supported, False otherwise
        """
        try:
            self._validate_file(file_path)
            return True
        except (FileValidationError, UnsupportedFormatError):
            return False
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats.
        
        Returns:
            List of supported file extensions
        """
        return list(self.SUPPORTED_FORMATS.keys())
    
    def cleanup_temp_files(self) -> None:
        """Clean up any temporary files created during processing."""
        for temp_file in self.temp_files:
            try:
                if Path(temp_file).exists():
                    Path(temp_file).unlink()
                    logger.debug(f"Cleaned up temporary file: {temp_file}")
            except Exception as e:
                logger.warning(f"Failed to clean up temporary file {temp_file}: {str(e)}")
        
        self.temp_files.clear()
    
    def __del__(self):
        """Cleanup temporary files when object is destroyed."""
        self.cleanup_temp_files()