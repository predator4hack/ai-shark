import os
import io
import tempfile
from pathlib import Path
from typing import List
from PIL import Image
from pptx import Presentation
import fitz  # PyMuPDF

class FileConverter:
    """Handles file format conversions"""
    
    @staticmethod
    def ppt_to_images(ppt_path: str) -> List[Image.Image]:
        """
        Extract images directly from PPT slides
        
        Args:
            ppt_path: Path to PowerPoint file
            
        Returns:
            List of PIL Images representing slides
        """
        try:
            presentation = Presentation(ppt_path)
            images = []
            
            for i, slide in enumerate(presentation.slides):
                # Create a temporary file to save slide as image
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                    temp_path = temp_file.name
                
                try:
                    # For now, we'll convert PPT to PDF first and then extract images
                    # This is a fallback approach since direct slide-to-image conversion
                    # requires additional dependencies
                    pdf_path = FileConverter.ppt_to_pdf(ppt_path)
                    if pdf_path:
                        return FileConverter.pdf_to_images(pdf_path)
                    else:
                        # If PDF conversion fails, create placeholder images
                        placeholder = Image.new('RGB', (800, 600), color='white')
                        images.append(placeholder)
                        
                except Exception as e:
                    print(f"Error processing slide {i}: {e}")
                    # Create placeholder image for failed slides
                    placeholder = Image.new('RGB', (800, 600), color='white')
                    images.append(placeholder)
                finally:
                    # Clean up temp file
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
            
            return images
            
        except Exception as e:
            print(f"Error converting PPT to images: {e}")
            return []
    
    @staticmethod
    def ppt_to_pdf(ppt_path: str) -> str:
        """
        Convert PPT to PDF using available system tools
        
        Args:
            ppt_path: Path to PowerPoint file
            
        Returns:
            Path to created PDF file, or None if conversion failed
        """
        try:
            # Create temporary PDF file
            temp_pdf = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
            pdf_path = temp_pdf.name
            temp_pdf.close()
            
            # Try using LibreOffice for conversion (if available)
            import subprocess
            try:
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', os.path.dirname(pdf_path), ppt_path
                ], capture_output=True, text=True, timeout=60)
                
                if result.returncode == 0:
                    # LibreOffice creates PDF with same name as input file
                    generated_pdf = os.path.join(
                        os.path.dirname(pdf_path),
                        Path(ppt_path).stem + '.pdf'
                    )
                    if os.path.exists(generated_pdf):
                        # Move to our temp location
                        os.rename(generated_pdf, pdf_path)
                        return pdf_path
                        
            except (subprocess.TimeoutExpired, FileNotFoundError):
                pass
            
            # If LibreOffice is not available, try alternative approach
            # For now, we'll use a simple conversion method
            print("LibreOffice not available, using fallback PPT processing")
            
            # Alternative: Extract text and create a simple PDF
            # This is a simplified approach for demonstration
            presentation = Presentation(ppt_path)
            
            # For now, return None to indicate PDF conversion is not available
            # The calling code will handle PPT directly
            os.unlink(pdf_path)
            return None
            
        except Exception as e:
            print(f"Error converting PPT to PDF: {e}")
            if 'pdf_path' in locals() and os.path.exists(pdf_path):
                os.unlink(pdf_path)
            return None
    
    @staticmethod
    def pdf_to_images(pdf_path: str) -> List[Image.Image]:
        """
        Convert PDF to list of images (reused from existing pipeline)
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            List of PIL Images representing pages
        """
        try:
            doc = fitz.open(pdf_path)
            images = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=150)
                img_data = pix.tobytes("png")
                image = Image.open(io.BytesIO(img_data))
                images.append(image)
            
            doc.close()
            return images
            
        except Exception as e:
            print(f"Error converting PDF to images: {e}")
            return []
    
    @staticmethod
    def extract_text_from_ppt(ppt_path: str) -> str:
        """
        Extract text content from PowerPoint file
        
        Args:
            ppt_path: Path to PowerPoint file
            
        Returns:
            Extracted text content
        """
        try:
            presentation = Presentation(ppt_path)
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.extend(slide_text)
                    text_content.append("")  # Empty line between slides
            
            return "\n".join(text_content)
            
        except Exception as e:
            print(f"Error extracting text from PPT: {e}")
            return ""